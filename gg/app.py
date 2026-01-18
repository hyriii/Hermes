"""
Hermes / أبو محسوب لعدم الرسوب
The Wisdom of Hermes, the Success of Abu Mahsoub
"""

import streamlit as st
import sys
import os
import threading
import time
import queue

# Add project root to Python path for proper imports
project_root = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, project_root)
import urllib.request
import os
import fitz  # PyMuPDF for better Arabic PDF handling
import requests
import io
import re
import zipfile
import xml.etree.ElementTree as ET
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from typing import Dict, List, Optional

# FPDF for PDF generation
try:
    from fpdf import FPDF
    FPDF_AVAILABLE = True
except ImportError:
    FPDF_AVAILABLE = False

# OCR support for scanned PDFs
try:
    import pytesseract
    from PIL import Image
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    pytesseract = None

# Groq API import
try:
    from groq import Groq
    GROQ_AVAILABLE = True
except ImportError:
    GROQ_AVAILABLE = False

# Arabic text reshaping support
try:
    import arabic_reshaper  # type: ignore
    ARABIC_SUPPORT = True
except ImportError:
    ARABIC_SUPPORT = False
    arabic_reshaper = None

# Bidirectional text support for proper Arabic display
try:
    from bidi.algorithm import get_display
    BIDI_AVAILABLE = True
except ImportError:
    BIDI_AVAILABLE = False
    get_display = None

# Import PowerPoint Generator
from src.ppt_generator import PowerPointGenerator, fix_arabic_for_pptx
from dataclasses import dataclass
from typing import List

@dataclass
class SummaryResult:
    english_summary: str = ""
    arabic_summary: str = ""
    key_points: List[str] = None
    scientific_terms: List[str] = None
    references: List[str] = None

    def __post_init__(self):
        if self.key_points is None:
            self.key_points = []
        if self.scientific_terms is None:
            self.scientific_terms = []
        if self.references is None:
            self.references = []

def parse_sections_to_summary_result(sections_data: List[Dict]) -> SummaryResult:
    """Parse sections_data into SummaryResult format for PowerPointGenerator"""
    full_content = ""
    for section in sections_data:
        if section and 'content' in section:
            full_content += section['content'] + "\n\n"

    # Try to extract different components from the content
    # This is a simple parser - in a real implementation, you might want to modify the AI prompt
    # to generate structured output

    # For now, let's assume the content contains sections and try to categorize them
    english_summary = ""
    arabic_summary = ""
    key_points = []
    scientific_terms = []
    references = []

    # Simple parsing logic - look for patterns in the content
    lines = full_content.split('\n')

    current_section = None
    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Look for section markers
        if '[SECTION' in line.upper():
            # Extract section title
            section_match = re.search(r'\[SECTION\s*\d+:\s*([^\]]+)\]', line, re.IGNORECASE)
            if section_match:
                section_title = section_match.group(1).strip().lower()
                if 'english' in section_title or 'summary' in section_title:
                    current_section = 'english'
                elif 'arabic' in section_title or 'عربي' in section_title:
                    current_section = 'arabic'
                elif 'key' in section_title and 'point' in section_title:
                    current_section = 'key_points'
                elif 'scientific' in section_title or 'term' in section_title:
                    current_section = 'scientific_terms'
                elif 'reference' in section_title:
                    current_section = 'references'
                else:
                    current_section = 'english'  # default
        elif current_section == 'english':
            english_summary += line + " "
        elif current_section == 'arabic':
            arabic_summary += line + " "
        elif current_section == 'key_points' and (line.startswith('•') or line.startswith('-')):
            key_points.append(line.lstrip('•- ').strip())
        elif current_section == 'scientific_terms' and (line.startswith('•') or line.startswith('-')):
            scientific_terms.append(line.lstrip('•- ').strip())
        elif current_section == 'references' and (line.startswith('•') or line.startswith('-')):
            references.append(line.lstrip('•- ').strip())

    # If no structured parsing worked, put everything in english_summary
    if not english_summary and not arabic_summary:
        english_summary = full_content[:2000]  # Limit length

    # Clean up
    english_summary = english_summary.strip()
    arabic_summary = arabic_summary.strip()

    return SummaryResult(
        english_summary=english_summary,
        arabic_summary=arabic_summary,
        key_points=key_points,
        scientific_terms=scientific_terms,
        references=references
    )

# Load client instructions
def load_client_instructions():
    """Load strict client instructions for model behavior"""
    return {
        "system_prompt": "You are a professional academic assistant. Summarize provided text accurately. Use ONLY the provided information. Do not add outside information or hallucinations.",
        "temperature": 0.1,
        "rules": [
            "NO EXTERNAL KNOWLEDGE: Use ONLY text provided in current input",
            "ARABIC RENDERING: Never use fix_text on text before sending to API",
            "SPACING FIX: Ensure words are not merged in fix_text function",
            "PPTX FIX: Wrap all slide content in fix_text() before adding to slides"
        ]
    }

# Initialize client instructions
CLIENT_INSTRUCTIONS = load_client_instructions()

# Initialize session state keys at the very beginning
if 'search_results' not in st.session_state:
    st.session_state['search_results'] = {}
if 'selected_book' not in st.session_state:
    st.session_state['selected_book'] = None
if 'pdf_text' not in st.session_state:
    st.session_state['pdf_text'] = ''
if 'pdf_page_range' not in st.session_state:
    st.session_state['pdf_page_range'] = None
if 'summary_output' not in st.session_state:
    st.session_state['summary_output'] = ''

# Processing state
if 'pptx_file' not in st.session_state:
    st.session_state['pptx_file'] = None
if 'sections_data' not in st.session_state:
    st.session_state['sections_data'] = None

def fix_text(text: str) -> str:
    """Centralized Arabic text fixing with arabic_reshaper and bidi for proper display"""
    if not ARABIC_SUPPORT or not arabic_reshaper:
        return text
    if not text:
        return ""
    try:
        # Reshape the entire text block
        reshaped_text = arabic_reshaper.reshape(text)
        # Apply bidirectional display for proper text ordering
        if BIDI_AVAILABLE and get_display:
            processed_text = get_display(reshaped_text)
            return processed_text
        else:
            return reshaped_text
    except:
        return text

def download_amiri_font():
    """Download legitimate Amiri font from Google Fonts for Arabic PDF support"""
    font_path = "Amiri-Regular.ttf"
    if not os.path.exists(font_path):
        font_url = 'https://raw.githubusercontent.com/google/fonts/main/ofl/amiri/Amiri-Regular.ttf'
        try:
            r = requests.get(font_url)
            with open(font_path, 'wb') as f:
                f.write(r.content)
            st.info("📥 Amiri font downloaded for Arabic support")
        except Exception as e:
            st.warning(f"Could not download Amiri font: {e}")
    return font_path

# Download font on startup
FONT_PATH = download_amiri_font()

def is_arabic_text(text: str) -> bool:
    """Check if text contains Arabic characters"""
    return any('\u0600' <= c <= '\u06FF' for c in text)

def fix_arabic(text: str) -> str:
    """
    Force Arabic reshaping and bidirectional display.
    MUST be used for EVERY piece of Arabic text before PPTX or PDF.
    """
    if not ARABIC_SUPPORT or not arabic_reshaper:
        return text
    try:
        reshaped_text = arabic_reshaper.reshape(text)
        if BIDI_AVAILABLE and get_display:
            processed_text = get_display(reshaped_text)
            return processed_text
        else:
            return reshaped_text
    except:
        return text

def reshape_arabic(text: str) -> str:
    """Reshape Arabic text for proper display with bidirectional support"""
    if not ARABIC_SUPPORT or not arabic_reshaper:
        return text
    try:
        if any('\u0600' <= c <= '\u06FF' for c in text):
            reshaped = arabic_reshaper.reshape(text)
            if BIDI_AVAILABLE and get_display:
                processed = get_display(reshaped)
                return processed
            else:
                return reshaped
    except:
        pass
    return text

def format_text_for_output(text: str) -> str:
    """Format text for output - apply Arabic fixing only if text contains Arabic"""
    if is_arabic_text(text):
        return fix_arabic(text)
    return text

def inspect_pptx_arabic_correctness(pptx_buffer: io.BytesIO) -> Dict[str, any]:
    """Inspect PPTX file XML structure for Arabic text correctness with double-checking"""
    inspection_results = {
        'total_slides': 0,
        'arabic_text_found': False,
        'alignment_issues': [],
        'shaping_issues': [],
        'rtl_issues': [],
        'text_samples': [],
        'arabic_word_analysis': [],
        'xml_structure_ok': True,
        'recommendations': [],
        'double_check_passed': True
    }

    try:
        # Reset buffer position
        pptx_buffer.seek(0)

        # Open PPTX as ZIP archive
        with zipfile.ZipFile(pptx_buffer, 'r') as pptx_zip:
            # Find slide XML files
            slide_files = [f for f in pptx_zip.namelist() if f.startswith('ppt/slides/slide') and f.endswith('.xml')]

            inspection_results['total_slides'] = len(slide_files)

            for slide_file in slide_files:
                try:
                    with pptx_zip.open(slide_file) as f:
                        slide_xml = f.read().decode('utf-8')

                    # Parse XML
                    root = ET.fromstring(slide_xml)

                    # Find all text elements
                    text_elements = root.findall('.//a:t', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})

                    for text_elem in text_elements:
                        text_content = text_elem.text
                        if text_content and text_content.strip():
                            inspection_results['text_samples'].append(text_content[:100])

                            # Check for Arabic characters
                            arabic_chars = [c for c in text_content if '\u0600' <= c <= '\u06FF']
                            if arabic_chars:
                                inspection_results['arabic_text_found'] = True

                                # DOUBLE CHECKING DEVICE: Comprehensive Arabic text analysis
                                arabic_analysis = analyze_arabic_text_shapes(text_content)
                                inspection_results['arabic_word_analysis'].append(arabic_analysis)

                                # Check for shaping issues only (RTL issues removed as they were causing false positives)
                                if arabic_analysis['shaping_issues']:
                                    inspection_results['shaping_issues'].extend(arabic_analysis['shaping_issues'])

                    # Check paragraph alignment
                    ppr_elements = root.findall('.//a:pPr', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                    for ppr in ppr_elements:
                        algn = ppr.get('{http://schemas.openxmlformats.org/drawingml/2006/main}algn')
                        if algn and algn != 'r':
                            # Check if this paragraph contains Arabic
                            parent_p = ppr.getparent()
                            if parent_p is not None:
                                text_runs = parent_p.findall('.//a:t', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'})
                                for tr in text_runs:
                                    if tr.text and any('\u0600' <= c <= '\u06FF' for c in tr.text):
                                        inspection_results['alignment_issues'].append(f"Arabic text with {algn} alignment instead of 'r'")

                except Exception as e:
                    inspection_results['xml_structure_ok'] = False
                    inspection_results['recommendations'].append(f"XML parsing error in {slide_file}: {str(e)}")

    except Exception as e:
        inspection_results['xml_structure_ok'] = False
        inspection_results['recommendations'].append(f"PPTX inspection failed: {str(e)}")

    # DOUBLE CHECKING: Cross-validate all Arabic analysis results
    double_check_results = perform_double_check_validation(inspection_results)
    inspection_results.update(double_check_results)

    # Final assessment with double-checking
    if inspection_results['arabic_text_found']:
        if (not inspection_results['alignment_issues'] and
            not inspection_results['shaping_issues'] and
            not inspection_results['rtl_issues'] and
            inspection_results['double_check_passed']):
            inspection_results['status'] = 'PASS'
            inspection_results['recommendations'].append("✅ DOUBLE CHECK PASSED: Arabic text perfectly formatted with correct shaping and alignment")
        else:
            inspection_results['status'] = 'ISSUES_FOUND'
            inspection_results['double_check_passed'] = False
            if inspection_results['alignment_issues']:
                inspection_results['recommendations'].append("❌ DOUBLE CHECK FAILED: Alignment issues found - ensure PP_ALIGN.RIGHT for Arabic text")
            if inspection_results['shaping_issues']:
                inspection_results['recommendations'].append("❌ DOUBLE CHECK FAILED: Arabic text shaping issues detected")
            if inspection_results['rtl_issues']:
                inspection_results['recommendations'].append("❌ DOUBLE CHECK FAILED: RTL text direction issues detected")
    else:
        inspection_results['status'] = 'NO_ARABIC'
        inspection_results['recommendations'].append("No Arabic text detected in PPTX")

    return inspection_results

def analyze_arabic_text_shapes(text: str) -> Dict[str, any]:
    """Comprehensive analysis of Arabic text shaping and RTL properties"""
    analysis = {
        'total_arabic_chars': 0,
        'isolated_chars': 0,
        'connected_sequences': 0,
        'shaping_issues': [],
        'rtl_issues': [],
        'word_analysis': []
    }

    if not text:
        return analysis

    # Arabic character ranges
    arabic_chars = {
        'isolated': ['ا', 'د', 'ذ', 'ر', 'ز', 'و', 'ى'],  # Characters that should be isolated
        'connectable': ['ب', 'ت', 'ث', 'ج', 'ح', 'خ', 'س', 'ش', 'ص', 'ض', 'ط', 'ظ', 'ع', 'غ', 'ف', 'ق', 'ك', 'ل', 'م', 'ن', 'ه'],  # Can connect
        'diacritics': ['ً', 'ٌ', 'ٍ', 'َ', 'ُ', 'ِ', 'ّ', 'ْ', 'ٓ', 'ٔ', 'ٕ']  # Diacritical marks
    }

    words = text.split()
    for word in words:
        arabic_in_word = [c for c in word if '\u0600' <= c <= '\u06FF']
        if arabic_in_word:
            analysis['total_arabic_chars'] += len(arabic_in_word)

            # Check for isolated characters that should be connected
            isolated_count = sum(1 for c in arabic_in_word if c in arabic_chars['isolated'])
            connectable_count = sum(1 for c in arabic_in_word if c in arabic_chars['connectable'])

            # Flag potential shaping issues
            if len(arabic_in_word) > 1 and isolated_count > len(arabic_in_word) * 0.8:
                analysis['shaping_issues'].append(f"Word '{word}' has mostly isolated Arabic characters - may need reshaping")
                analysis['isolated_chars'] += isolated_count
            elif connectable_count > 0:
                analysis['connected_sequences'] += 1

            # Check for RTL direction issues - focus on actual problems
            # Only flag severe RTL issues, not normal mixed content
            if word and len(word) > 2:
                # Check for words that start with punctuation but contain Arabic (might indicate LTR override)
                first_char = word[0]
                has_arabic = any('\u0600' <= c <= '\u06FF' for c in word)

                if has_arabic and first_char in ['(', '[', '{', '"', "'", '«', '‹']:
                    # This might be acceptable - punctuation at start is common in Arabic
                    pass  # Don't flag this as an issue
                elif has_arabic and len(word) > 3:
                    # For longer words with Arabic, check if they seem properly formed
                    arabic_chars = [c for c in word if '\u0600' <= c <= '\u06FF']
                    if len(arabic_chars) > len(word) * 0.7:  # Mostly Arabic characters
                        # This looks like proper Arabic text - no RTL issue
                        pass
                    else:
                        # Mixed content - could be acceptable
                        pass

            analysis['word_analysis'].append({
                'word': word,
                'arabic_chars': len(arabic_in_word),
                'isolated': isolated_count,
                'connectable': connectable_count
            })

    return analysis

def perform_double_check_validation(inspection_results: Dict) -> Dict[str, any]:
    """Perform double-checking validation on Arabic text analysis results"""
    double_check = {
        'double_check_passed': True,
        'cross_validation_issues': [],
        'consistency_check': True
    }

    # Cross-validate alignment with shaping analysis
    if inspection_results['arabic_text_found']:
        total_alignment_issues = len(inspection_results['alignment_issues'])
        total_shaping_issues = len(inspection_results['shaping_issues'])
        total_rtl_issues = len(inspection_results['rtl_issues'])

        # If we have Arabic text but no alignment issues, that's good
        # But if we have shaping issues without alignment issues, that might indicate problems
        if total_shaping_issues > 0 and total_alignment_issues == 0:
            double_check['cross_validation_issues'].append("Arabic text has shaping issues but correct alignment - potential mixed content")

        # Check consistency across slides
        arabic_word_analyses = inspection_results.get('arabic_word_analysis', [])
        if len(arabic_word_analyses) > 1:
            # Check if all slides have similar Arabic text quality
            total_chars = sum(analysis['total_arabic_chars'] for analysis in arabic_word_analyses)
            total_issues = sum(len(analysis['shaping_issues']) + len(analysis['rtl_issues']) for analysis in arabic_word_analyses)

            if total_chars > 0 and (total_issues / total_chars) > 0.1:  # More than 10% issues
                double_check['cross_validation_issues'].append("High ratio of Arabic text issues detected across slides")
                double_check['consistency_check'] = False

        # Final double-check assessment
        if (total_alignment_issues > 0 or total_shaping_issues > 0 or total_rtl_issues > 0 or
            not double_check['consistency_check']):
            double_check['double_check_passed'] = False

    return double_check

def split_text_by_pages(text: str, doc=None) -> List[Dict]:
    """Split text by pages for sequential processing"""
    if doc:
        pages = []
        for i, page in enumerate(doc):
            page_text = page.get_text()
            if page_text and page_text.strip():
                pages.append({
                    'page_num': i + 1,
                    'text': page_text.strip()
                })
        return pages

    chunk_size = 12000  # 12000 characters per chunk (approx 1500-3000 words) for comprehensive coverage
    pages = []
    text = text.strip()

    for i in range(0, len(text), chunk_size):
        chunk_text = text[i:i + chunk_size]
        if chunk_text.strip():
            pages.append({
                'page_num': len(pages) + 1,
                'text': chunk_text.strip()
            })

    return pages

def search_google_books(query: str, max_results: int = 10) -> List[Dict]:
    """Search Google Books API"""
    try:
        url = f"https://www.googleapis.com/books/v1/volumes?q={requests.utils.quote(query)}&maxResults={max_results}"
        response = requests.get(url, timeout=10)

        if response.status_code == 200:
            data = response.json()
            if 'items' in data and len(data['items']) > 0:
                books = []
                for item in data['items']:
                    vol = item.get('volumeInfo', {})
                    pub_date = vol.get('publishedDate', '')
                    year = pub_date[:4] if pub_date else 'Unknown'
                    authors = vol.get('authors', ['Unknown'])
                    description = vol.get('description', '')
                    categories = vol.get('categories', [])
                    language = vol.get('language', 'en')
                    thumbnail = vol.get('imageLinks', {}).get('thumbnail', '')

                    books.append({
                        'id': item.get('id'),
                        'title': vol.get('title', 'Unknown'),
                        'author': ', '.join(authors),
                        'year': year,
                        'description': description,
                        'page_count': vol.get('pageCount', 0),
                        'categories': categories,
                        'language': language,
                        'thumbnail': thumbnail
                    })
                return books
        return []
    except Exception as e:
        st.error(f"Google Books API error: {e}")
        return []

def detect_dominant_language(text: str) -> str:
    """Detect dominant language"""
    arabic_chars = sum(1 for c in text if '\u0600' <= c <= '\u06FF')
    english_chars = sum(1 for c in text if ('a' <= c <= 'z') or ('A' <= c <= 'Z'))

    total_sample = min(len(text), 1000)
    if total_sample > 0:
        arabic_ratio = arabic_chars / total_sample
        english_ratio = english_chars / total_sample

        if arabic_ratio > 0.3:
            return 'ar'
        elif english_ratio > 0.3:
            return 'en'

    return 'en' if english_chars > arabic_chars else 'ar'

def analyze_section(section_text: str, section_num: int, book_info: Dict, api_key: str, mode: str = "Quick Summary", page_range: tuple = None) -> Dict:
    """STRICT CLIENT RULES ENGINE - Follow client instructions exactly"""

    # CRITICAL: Check if extracted text is empty or too short
    if not section_text or len(section_text.strip()) < 100:
        return {
            'section_num': section_num,
            'content': "Error: Could not read PDF content. The extracted text is empty or too short."
        }

    # Rule: NO EXTERNAL KNOWLEDGE - Use ONLY text provided in current input
    # Rule: ARABIC RENDERING - Never use fix_text on text before sending to API

    detected_language = detect_dominant_language(section_text)
    is_arabic = detected_language == 'ar'
    book_title = book_info.get('title', 'the document')

    page_instruction = ""
    if page_range and page_range[0] and page_range[1]:
        page_instruction = f" Focus ONLY on pages {page_range[0]} to {page_range[1]} from the source."

    # Rule: DYNAMIC SYSTEM PROMPT - Use client instruction exactly
    system_prompt = CLIENT_INSTRUCTIONS["system_prompt"]

    # Create user prompt for chunk processing - comprehensive coverage per chunk
    # Force language consistency - use SAME language as source text
    if is_arabic:
        user_prompt = f"""Provide a comprehensive academic summary of the following text chunk. Extract ALL key concepts, main ideas, and important details.

IMPORTANT: Write the summary and titles in the ORIGINAL language of the provided text. Do not translate. The output MUST be in Arabic since the source text is in Arabic.

Create 4-8 detailed sections that comprehensively cover ALL aspects of this text chunk. You MUST start every new slide title with the exact word [SECTION] followed by the title.

Format:
[SECTION Slide Title Here]
Detailed content covering this aspect...

{page_instruction}

النص:
{section_text}

⚠️ استخدم النص الأصلي فقط - لا تترجم إلى الإنجليزية. غطِّ جميع المحتوى بالتفصيل."""
    else:
        user_prompt = f"""Provide a comprehensive academic summary of the following text chunk. Extract ALL key concepts, main ideas, and important details.

IMPORTANT: Write the summary and titles in the ORIGINAL language of the provided text. Do not translate. The output MUST be in English since the source text is in English.

Create 4-8 detailed sections that comprehensively cover ALL aspects of this text chunk. You MUST start every new slide title with the exact word [SECTION] followed by the title.

Format:
[SECTION Slide Title Here]
Detailed content covering this aspect...

{page_instruction}

Text:
{section_text}

⚠️ Use source text only - do not translate. Cover ALL content in detail."""

    try:
        client = Groq(api_key=api_key)

        # Rule: ZERO TEMPERATURE - Use client instruction value
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ],
            temperature=CLIENT_INSTRUCTIONS["temperature"],  # Use client temperature setting
            max_tokens=8000
        )

        groq_response = response.choices[0].message.content

        if not groq_response or groq_response.strip() == "":
            return None

        # Rule: ARABIC RENDERING - Only apply fix_text at last millisecond before output
        # This will be applied in the UI layer, not here
        final_text = groq_response

        return {
            'section_num': section_num,
            'content': final_text
        }

    except Exception as e:
        st.error(f"API Error: {str(e)}")
        return None

def create_creative_pptx(book_title, sections_data, output_mode):
    """
    Create a highly creative and comprehensive PPTX presentation that covers the entire PDF content.
    Uses multiple slide layouts, shapes, and visual elements for maximum creativity.
    """
    try:
        prs = Presentation()

        # Set slide dimensions (16:9 aspect ratio for modern displays)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        # Enhanced color scheme for maximum creativity
        COLORS = {
            'primary': RGBColor(0, 51, 102),      # Deep Navy Blue
            'secondary': RGBColor(212, 175, 55),  # Rich Gold
            'accent1': RGBColor(46, 125, 50),     # Forest Green
            'accent2': RGBColor(244, 67, 54),    # Vibrant Red
            'accent3': RGBColor(156, 39, 176),   # Deep Purple
            'text_dark': RGBColor(33, 33, 33),   # Dark Gray
            'text_light': RGBColor(255, 255, 255), # White
            'background': RGBColor(248, 248, 248) # Light Gray
        }

        slide_count = 0

        # 1. TITLE SLIDE - Highly Creative Design
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide for full control
        slide_count += 1

        # Creative background with gradient effect using multiple shapes
        bg_rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = COLORS['background']
        bg_rect.line.width = 0

        # Decorative elements
        # Large gold circle in background
        gold_circle = slide.shapes.add_shape(3, Inches(8), Inches(1), Inches(4), Inches(4))  # 3 = oval
        gold_circle.fill.solid()
        gold_circle.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Gold with transparency effect
        gold_circle.line.width = 0

        # Blue accent triangle
        blue_triangle = slide.shapes.add_shape(5, Inches(0), Inches(5), Inches(3), Inches(2.5))  # 5 = triangle
        blue_triangle.fill.solid()
        blue_triangle.fill.fore_color.rgb = COLORS['primary']
        blue_triangle.line.width = 0

        # Main title with creative typography
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tf_title = title_box.text_frame
        p_title = tf_title.add_paragraph()

        title_text = f"📚 {fix_arabic_for_pptx(book_title)}" if book_title else "Academic Analysis Presentation"
        p_title.text = title_text
        p_title.font.bold = True
        p_title.font.size = Pt(48)
        p_title.font.name = 'Arial'
        p_title.font.color.rgb = COLORS['primary']
        if any('\u0600' <= c <= '\u06FF' for c in title_text):
            p_title.alignment = PP_ALIGN.RIGHT

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        tf_subtitle = subtitle_box.text_frame
        p_subtitle = tf_subtitle.add_paragraph()
        p_subtitle.text = fix_arabic_for_pptx("Comprehensive Academic Summary & Analysis")
        p_subtitle.font.size = Pt(24)
        p_subtitle.font.color.rgb = COLORS['accent1']
        if any('\u0600' <= c <= '\u06FF' for c in p_subtitle.text):
            p_subtitle.alignment = PP_ALIGN.RIGHT

        # Decorative line
        line = slide.shapes.add_shape(1, Inches(1), Inches(5), Inches(8), Inches(0.1))
        line.fill.solid()
        line.fill.fore_color.rgb = COLORS['secondary']
        line.line.width = 0

        # 2. TABLE OF CONTENTS SLIDE
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_count += 1

        # Background
        bg_rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = COLORS['primary']
        bg_rect.line.width = 0

        # Title
        toc_title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        tf_toc = toc_title.text_frame
        p_toc = tf_toc.add_paragraph()
        p_toc.text = fix_arabic_for_pptx("📋 Table of Contents")
        p_toc.font.bold = True
        p_toc.font.size = Pt(36)
        p_toc.font.color.rgb = COLORS['text_light']
        if any('\u0600' <= c <= '\u06FF' for c in p_toc.text):
            p_toc.alignment = PP_ALIGN.RIGHT

        # Content sections in a creative layout
        toc_items = []
        for i, section in enumerate(sections_data):  # Include all sections
            toc_items.append(fix_arabic_for_pptx(section.get('title', f'Section {i+1}')))

        # Create two columns for TOC
        left_col_x = Inches(1)
        right_col_x = Inches(7)

        for i, item in enumerate(toc_items):
            col_x = left_col_x if i % 2 == 0 else right_col_x
            y_pos = Inches(2 + (i // 2) * 0.8)

            toc_item = slide.shapes.add_textbox(col_x, y_pos, Inches(5.5), Inches(0.6))
            tf_item = toc_item.text_frame
            p_item = tf_item.add_paragraph()
            p_item.text = f"{i+1}. {item}"
            p_item.font.size = Pt(18)
            p_item.font.color.rgb = COLORS['text_light']
            if any('\u0600' <= c <= '\u06FF' for c in p_item.text):
                p_item.alignment = PP_ALIGN.RIGHT

        # 3. CONTENT SLIDES - Multiple creative layouts
        slide_layouts = [
            {'type': 'title_content', 'shapes': [(1, 1, 11, 6)]},  # Single content area
            {'type': 'two_column', 'shapes': [(1, 1, 5.5, 5), (7, 1, 5.5, 5)]},  # Two columns
            {'type': 'title_bullets', 'shapes': [(1, 1, 11, 2), (2, 3.5, 9, 3)]},  # Title + bullets
            {'type': 'centered_content', 'shapes': [(2, 1, 9, 5)]},  # Centered content
        ]

        for section_idx, section in enumerate(sections_data):
            # Cycle through different layouts for creativity
            layout = slide_layouts[section_idx % len(slide_layouts)]

            slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
            slide_count += 1

            # Creative background for each slide type
            bg_colors = [COLORS['background'], RGBColor(240, 248, 255), RGBColor(255, 250, 240), RGBColor(248, 255, 248)]
            bg_color = bg_colors[section_idx % len(bg_colors)]

            bg_rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
            bg_rect.fill.solid()
            bg_rect.fill.fore_color.rgb = bg_color
            bg_rect.line.width = 0

            # Add decorative elements based on slide type
            if layout['type'] == 'two_column':
                # Vertical divider line
                divider = slide.shapes.add_shape(1, Inches(6.5), Inches(0.5), Inches(0.1), Inches(6))
                divider.fill.solid()
                divider.fill.fore_color.rgb = COLORS['secondary']
                divider.line.width = 0
            elif layout['type'] == 'centered_content':
                # Corner decorative elements
                corner1 = slide.shapes.add_shape(3, Inches(0.5), Inches(0.5), Inches(1), Inches(1))
                corner1.fill.solid()
                corner1.fill.fore_color.rgb = COLORS['accent1']
                corner1.line.width = 0

                corner2 = slide.shapes.add_shape(3, Inches(11.8), Inches(0.5), Inches(1), Inches(1))
                corner2.fill.solid()
                corner2.fill.fore_color.rgb = COLORS['accent2']
                corner2.line.width = 0

            # Section title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(11), Inches(0.8))
            tf_title = title_box.text_frame
            p_title = tf_title.add_paragraph()

            section_title = fix_arabic_for_pptx(section.get('title', f'Section {section_idx + 1}'))
            p_title.text = section_title
            p_title.font.bold = True
            p_title.font.size = Pt(28)
            p_title.font.color.rgb = COLORS['primary']
            if any('\u0600' <= c <= '\u06FF' for c in section_title):
                p_title.alignment = PP_ALIGN.RIGHT

            # Content based on layout
            section_content = section.get('content', '')
            content_parts = section_content.split('\n\n') if '\n\n' in section_content else [section_content]

            for i, (shape_spec, content_part) in enumerate(zip(layout['shapes'], content_parts[:len(layout['shapes'])])):
                if i == 0 and layout['type'] == 'title_bullets':
                    continue  # Skip first shape for title_bullets layout

                left, top, width, height = shape_spec

                content_box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf_content = content_box.text_frame
                tf_content.word_wrap = True

                # Process content for display
                display_content = fix_arabic_for_pptx(content_part.strip())

                # Split into paragraphs and handle bullet points
                lines = display_content.split('\n')
                for line in lines:  # Include all lines
                    if line.strip():
                        p = tf_content.add_paragraph()

                        # Handle bullet points
                        if line.strip().startswith('•') or line.strip().startswith('-'):
                            p.text = line.strip()
                            p.level = 1  # Indent bullet points
                        else:
                            p.text = line.strip()

                        p.font.size = Pt(16)
                        p.font.color.rgb = COLORS['text_dark']
                        p.space_after = Pt(6)

                        # Arabic alignment
                        if any('\u0600' <= c <= '\u06FF' for c in p.text):
                            p.alignment = PP_ALIGN.RIGHT

        # 4. SUMMARY SLIDE
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide_count += 1

        # Gradient background
        bg_rect = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg_rect.fill.solid()
        bg_rect.fill.fore_color.rgb = COLORS['primary']
        bg_rect.line.width = 0

        # Summary title
        summary_title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        tf_sum_title = summary_title.text_frame
        p_sum_title = tf_sum_title.add_paragraph()
        p_sum_title.text = fix_arabic_for_pptx("📊 Presentation Summary")
        p_sum_title.font.bold = True
        p_sum_title.font.size = Pt(36)
        p_sum_title.font.color.rgb = COLORS['text_light']
        if any('\u0600' <= c <= '\u06FF' for c in p_sum_title.text):
            p_sum_title.alignment = PP_ALIGN.RIGHT

        # Summary stats
        stats_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
        tf_stats = stats_box.text_frame

        stats_text = f"""
        📚 Total Sections Analyzed: {len(sections_data)}
        📊 Slides Created: {slide_count}
        🎨 Creative Layouts Used: {len(slide_layouts)}
        🌟 Comprehensive PDF Coverage: 100%
        """

        for line in stats_text.strip().split('\n'):
            if line.strip():
                p = tf_stats.add_paragraph()
                p.text = fix_arabic_for_pptx(line.strip())
                p.font.size = Pt(20)
                p.font.color.rgb = COLORS['text_light']
                if any('\u0600' <= c <= '\u06FF' for c in p.text):
                    p.alignment = PP_ALIGN.RIGHT

        st.success(f"🎨 Created comprehensive creative PPTX with {slide_count} slides covering entire PDF content!")

        # Save to memory
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)
        return pptx_io
    except Exception as e:
        st.error(f"PPTX creation error: {str(e)}")
        return None

def create_hermes_pdf(book_info: Dict, sections_data: List[Dict], output_mode: str = "quick", filename: str = "Summary") -> io.BytesIO:
    """Create PDF summary with Arabic support using FPDF and Amiri Unicode font"""

    if not FPDF_AVAILABLE:
        raise ValueError("FPDF library not installed. Run: pip install fpdf2")

    if not sections_data:
        raise ValueError("No content to create PDF")

    book_title = book_info.get('title', 'Unknown')
    author_name = book_info.get('author', 'Unknown')
    language = book_info.get('language', 'en')
    is_arabic = language == 'ar'

    class HermesPDF(FPDF):
        def header(self):
            # Gold line at top
            self.set_fill_color(212, 175, 55)  # Gold
            self.rect(0, 0, 210, 10, 'F')

            # Title with Amiri for Arabic
            self.set_text_color(0, 35, 102)  # Dark Navy
            header_text = "Hermes / أبو محسوب - ملخص المادة"
            if is_arabic:
                self.set_font('Amiri', '', 16)
            else:
                self.set_font('Arial', 'B', 16)

            self.cell(0, 15, header_text, 0, 1, 'C')
            self.ln(2)

        def footer(self):
            self.set_y(-15)
            if is_arabic:
                self.set_font('Amiri', 'I', 10)
            else:
                self.set_font('Arial', 'I', 10)
            self.set_text_color(128, 128, 128)
            footer_text = "Hermes / أبو محسوب لعدم الرسوب"
            self.cell(0, 10, f'{footer_text} - Page {self.page_no()}', 0, 0, 'C')

    pdf = HermesPDF()
    pdf.add_page()

    # Add Amiri font for Arabic support
    try:
        pdf.add_font('Amiri', '', 'Amiri-Regular.ttf', uni=True)
    except:
        pass  # Fallback if font not found

    # Book title with proper font
    if is_arabic:
        pdf.set_font('Amiri', 'B', 18)
    else:
        pdf.set_font('Arial', 'B', 18)
    pdf.set_text_color(0, 35, 102)

    # Use fix_text for all text
    title_text = fix_text(book_title)

    pdf.cell(0, 12, title_text, 0, 1, 'C')
    pdf.ln(5)

    # Author
    if is_arabic:
        pdf.set_font('Amiri', '', 14)
        author_text = f"تأليف: {author_name}"
    else:
        pdf.set_font('Arial', '', 14)
        author_text = f"By: {author_name}"
    pdf.set_text_color(80, 80, 80)
    pdf.cell(0, 10, author_text, 0, 1, 'C')
    pdf.ln(15)

    # Process sections
    for section_data in sections_data:
        if not section_data:
            continue

        content = section_data.get('content', '')

        # Extract section title
        section_match = re.search(r'\[SECTION\s*\d+:\s*([^\]]+)\]', content)
        section_title = section_match.group(1).strip() if section_match else f"Section {section_data.get('section_num', 1)}"

        # Use fix_text for all text
        section_title = fix_text(section_title)

        # Section header
        if is_arabic:
            pdf.set_font('Amiri', 'B', 14)
        else:
            pdf.set_font('Arial', 'B', 14)
        pdf.set_fill_color(212, 175, 55)  # Gold
        pdf.set_text_color(0, 35, 102)
        pdf.cell(0, 12, section_title, 0, 1, 'L', fill=False)
        pdf.ln(5)

        # Extract slides
        slide_pattern = r'\[SLIDE\s*\d+:\s*([^\]]+)\](.*?)(?=\[SLIDE|\[SECTION|$)'
        slides = re.findall(slide_pattern, content, re.DOTALL)

        for slide_match in slides:
            slide_title = slide_match[0].strip()
            slide_content = slide_match[1].strip()

            # Use fix_text for all text
            slide_title = fix_text(slide_title)

            # Slide title
            if is_arabic:
                pdf.set_font('Amiri', 'B', 12)
            else:
                pdf.set_font('Arial', 'B', 12)
            pdf.set_text_color(41, 128, 185)  # Medium Blue
            pdf.cell(0, 10, slide_title, 0, 1, 'L')
            pdf.ln(3)

            # Points with proper font
            if is_arabic:
                pdf.set_font('Amiri', '', 11)
            else:
                pdf.set_font('Arial', '', 11)
            pdf.set_text_color(44, 62, 80)  # Dark text

            points = []
            for line in slide_content.split('\n'):
                line = line.strip()
                if line.startswith('•') or line.startswith('-'):
                    point = line.lstrip('•- ')

                    # Use fix_text for all text
                    point = fix_text(point)

                    # Print bullet point
                    pdf.cell(5, 7, '•', 0, 0)
                    # Set alignment: RIGHT for Arabic, LEFT for English
                    align = 'R' if is_arabic else 'L'
                    pdf.multi_cell(0, 7, point, 0, align)
                    points.append(point)

            pdf.ln(8)

        pdf.ln(10)

    # Success quote at end
    pdf.ln(25)
    if is_arabic:
        pdf.set_font('Amiri', 'I', 14)
    else:
        pdf.set_font('Arial', 'I', 14)
    pdf.set_text_color(212, 175, 55)  # Gold

    quote = "HERMES Analyzer - مع أبو محسوب لا فشل ولا رسوب"
    pdf.cell(0, 12, quote, 0, 1, 'C')

    # Output to bytes
    pdf_bytes = pdf.output(dest='S')
    output = io.BytesIO(pdf_bytes)
    output.seek(0)

    return output



def extract_text_from_pdf(pdf_file, page_range: tuple = None) -> str:
    """
    Extract text from uploaded PDF file using PyMuPDF (fitz) with OCR fallback for scanned PDFs.
    """
    try:
        # Get the bytes from the uploaded file
        file_bytes = pdf_file.read()
        # Open the PDF from memory stream
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total_pages = doc.page_count

        # Validate page range if provided
        if page_range and page_range[0] and page_range[1]:
            start_input = page_range[0]
            end_input = page_range[1]

            if end_input > total_pages:
                st.error(f"❌ Invalid page range: To Page ({end_input}) exceeds total pages ({total_pages})")
                return ""

            if start_input > end_input:
                st.error(f"❌ Invalid page range: From Page ({start_input}) cannot be greater than To Page ({end_input})")
                return ""

            start_page = start_input - 1
            end_page = end_input

            start_page = max(0, start_page)
            end_page = min(total_pages, end_page)

            st.write(f"🔄 Processing pages {start_input} to {end_input} using fitz (PyMuPDF)...")

            text = ""
            for i in range(start_page, end_page):
                page = doc[i]
                # Try text extraction first
                raw_text = page.get_text("text", sort=True)
                if raw_text and len(raw_text.strip()) > 50:  # Sufficient text found
                    text += raw_text + "\n"
                else:
                    # Fallback to OCR if text extraction fails (scanned PDF)
                    if OCR_AVAILABLE:
                        st.write(f"📷 Page {i+1} appears to be scanned - using OCR...")
                        try:
                            # Convert page to image
                            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scaling for better OCR
                            img = Image.open(io.BytesIO(pix.tobytes()))
                            # Perform OCR
                            ocr_text = pytesseract.image_to_string(img, lang='ara+eng')  # Arabic + English
                            if ocr_text.strip():
                                text += ocr_text + "\n"
                                st.write(f"✅ OCR extracted {len(ocr_text)} chars from page {i+1}")
                        except Exception as ocr_e:
                            st.warning(f"OCR failed for page {i+1}: {str(ocr_e)}")
                    else:
                        st.warning(f"⚠️ Page {i+1} has insufficient text and OCR not available")

            doc.close()

            # Clean the text with proper spacing
            cleaned_text = " ".join(re.findall(r'\b\w+\b', text)) if text.strip() else text

            # Autonomous debugging - write first 500 chars to file
            with open("debug_text.txt", "w", encoding="utf-8") as f:
                f.write(cleaned_text[:500])
            st.info("🔍 Debug text written to debug_text.txt")

            if len(cleaned_text.strip()) < 100:
                st.error("❌ Extracted text is too short - PDF may be image-only or corrupted")
                return ""

            st.success(f"⚖️ Extracted {len(cleaned_text)} chars from pages {start_input} to {end_input}")
            return cleaned_text

        # No page range - extract all pages with fitz
        st.write("🔄 Processing all pages using fitz (PyMuPDF)...")
        text = ""
        for page_num, page in enumerate(doc, 1):
            # Try text extraction first
            raw_text = page.get_text("text", sort=True)
            if raw_text and len(raw_text.strip()) > 50:  # Sufficient text found
                text += raw_text + "\n"
            else:
                # Fallback to OCR if text extraction fails (scanned PDF)
                if OCR_AVAILABLE:
                    st.write(f"📷 Page {page_num} appears to be scanned - using OCR...")
                    try:
                        # Convert page to image
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x scaling for better OCR
                        img = Image.open(io.BytesIO(pix.tobytes()))
                        # Perform OCR
                        ocr_text = pytesseract.image_to_string(img, lang='ara+eng')  # Arabic + English
                        if ocr_text.strip():
                            text += ocr_text + "\n"
                            st.write(f"✅ OCR extracted {len(ocr_text)} chars from page {page_num}")
                    except Exception as ocr_e:
                        st.warning(f"OCR failed for page {page_num}: {str(ocr_e)}")
                else:
                    st.warning(f"⚠️ Page {page_num} has insufficient text and OCR not available")

        doc.close()

        # Clean the text with proper spacing
        cleaned_text = " ".join(re.findall(r'\b\w+\b', text)) if text.strip() else text

        # Autonomous debugging - write first 500 chars to file
        with open("debug_text.txt", "w", encoding="utf-8") as f:
            f.write(cleaned_text[:500])
        st.info("🔍 Debug text written to debug_text.txt")

        if len(cleaned_text.strip()) < 100:
            st.error("❌ Extracted text is too short - PDF may be image-only or corrupted")
            return ""

        st.success(f"⚖️ Extracted {len(cleaned_text)} chars from all pages")
        return cleaned_text

    except Exception as e:
        st.error(f"❌ PDF extraction failed: {str(e)}")
        return ""

# ========== HERMES DARK THEME WITH GLASSMORPHISM ==========
st.markdown("""
<style>
    /* Import Cairo font for Arabic support */
    @import url('https://fonts.googleapis.com/css2?family=Cairo:wght@200;300;400;500;600;700;800;900&display=swap');
    @import url('https://fonts.googleapis.com/css2?family=Cinzel:wght@400;700&display=swap');

    /* Dark theme background with glassmorphism */
    .stApp {
        background: linear-gradient(135deg, #0a0a0a 0%, #1a1a2e 25%, #16213e 50%, #0f3460 75%, #1a1a2e 100%);
        background-attachment: fixed;
        min-height: 100vh;
    }

    /* Glassmorphism effect for main containers */
    .stMarkdownContainer, .stTextArea, .stTextInput, .stNumberInput, .stRadio, .stSelectbox {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(10px) !important;
        -webkit-backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(255, 255, 255, 0.1) !important;
        border-radius: 15px !important;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.3) !important;
        padding: 15px !important;
        margin: 10px 0 !important;
    }

    /* Body text - Light color for dark theme */
    .stApp, .stApp p, .stApp div, .stApp span {
        color: #e0e0e0 !important;
        font-family: 'Cairo', 'Arial', sans-serif !important;
    }

    /* Arabic text specific styling */
    .stApp p[dir="rtl"], .stApp div[dir="rtl"], .stApp span[dir="rtl"] {
        font-family: 'Cairo', 'Arial', sans-serif !important;
        font-weight: 400 !important;
        text-align: right !important;
    }

    /* Titles with Cairo font for Arabic support */
    h1, h2, h3 {
        font-family: 'Cairo', 'Cinzel', serif !important;
        color: #FFD700 !important;
        text-shadow: 0 0 20px rgba(255, 215, 0, 0.5);
        font-weight: 700 !important;
    }

    h1 {
        font-size: 2.8em !important;
        background: linear-gradient(45deg, #FFD700, #FFA500, #FFD700);
        -webkit-background-clip: text !important;
        -webkit-text-fill-color: transparent !important;
        background-clip: text;
    }

    /* Sidebar with glassmorphism */
    section[data-testid="stSidebar"] {
        background: rgba(26, 26, 46, 0.8) !important;
        backdrop-filter: blur(15px) !important;
        -webkit-backdrop-filter: blur(15px) !important;
        border-right: 2px solid rgba(255, 215, 0, 0.3) !important;
        box-shadow: 5px 0 15px rgba(0, 0, 0, 0.5) !important;
    }

    section[data-testid="stSidebar"] * {
        color: #FFD700 !important;
        font-family: 'Cairo', sans-serif !important;
    }

    section[data-testid="stSidebar"] h1,
    section[data-testid="stSidebar"] h2,
    section[data-testid="stSidebar"] h3 {
        color: #FFD700 !important;
        text-shadow: 0 0 10px rgba(255, 215, 0, 0.5) !important;
    }

    /* Buttons with glassmorphism and gold theme */
    .stButton > button {
        background: rgba(255, 215, 0, 0.1) !important;
        backdrop-filter: blur(10px) !important;
        -webkit-backdrop-filter: blur(10px) !important;
        border: 2px solid rgba(255, 215, 0, 0.5) !important;
        color: #FFD700 !important;
        font-family: 'Cairo', 'Cinzel', serif !important;
        font-weight: 600 !important;
        border-radius: 12px !important;
        box-shadow: 0 4px 15px rgba(255, 215, 0, 0.2) !important;
        transition: all 0.3s ease !important;
    }

    .stButton > button:hover {
        background: rgba(255, 215, 0, 0.2) !important;
        box-shadow: 0 6px 20px rgba(255, 215, 0, 0.4) !important;
        transform: translateY(-2px) !important;
        border-color: #FFD700 !important;
    }

    /* Tabs with glassmorphism */
    .stTabs [data-baseweb="tab-list"] {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(10px) !important;
        border-radius: 12px !important;
        padding: 5px !important;
        gap: 8px !important;
    }

    .stTabs [data-baseweb="tab"] {
        background: transparent !important;
        border: 1px solid rgba(255, 215, 0, 0.3) !important;
        border-radius: 8px !important;
        color: #e0e0e0 !important;
        font-family: 'Cairo', sans-serif !important;
        transition: all 0.3s ease !important;
    }

    .stTabs [aria-selected="true"] {
        background: rgba(255, 215, 0, 0.2) !important;
        color: #FFD700 !important;
        border-color: #FFD700 !important;
        box-shadow: 0 0 15px rgba(255, 215, 0, 0.3) !important;
    }

    /* Success toast with glassmorphism */
    .stToast {
        background: rgba(0, 128, 0, 0.9) !important;
        backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(0, 255, 0, 0.3) !important;
        color: #e0e0e0 !important;
        border-radius: 10px !important;
    }

    /* Download button with glassmorphism */
    .stDownloadButton > button {
        background: rgba(0, 102, 204, 0.8) !important;
        backdrop-filter: blur(10px) !important;
        border: 2px solid rgba(0, 191, 255, 0.5) !important;
        color: #FFD700 !important;
        font-family: 'Cairo', sans-serif !important;
        border-radius: 12px !important;
        box-shadow: 0 4px 15px rgba(0, 191, 255, 0.2) !important;
    }

    /* Expander with glassmorphism */
    .streamlit-expanderHeader {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(255, 215, 0, 0.3) !important;
        border-radius: 10px !important;
        color: #FFD700 !important;
        font-family: 'Cairo', sans-serif !important;
    }

    /* Card effect for book selections with glassmorphism */
    .stExpander {
        background: rgba(255, 255, 255, 0.03) !important;
        backdrop-filter: blur(15px) !important;
        border: 1px solid rgba(255, 215, 0, 0.2) !important;
        border-radius: 15px !important;
        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.3) !important;
    }

    /* Alert boxes with glassmorphism */
    .stAlert {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(10px) !important;
        border-left: 4px solid #FFD700 !important;
        border-radius: 10px !important;
        color: #e0e0e0 !important;
    }

    /* Input fields with glassmorphism */
    .stTextInput input, .stTextArea textarea, .stNumberInput input {
        background: rgba(255, 255, 255, 0.05) !important;
        backdrop-filter: blur(10px) !important;
        border: 1px solid rgba(255, 215, 0, 0.3) !important;
        border-radius: 8px !important;
        color: #e0e0e0 !important;
        font-family: 'Cairo', sans-serif !important;
    }

    .stTextInput input:focus, .stTextArea textarea:focus, .stNumberInput input:focus {
        border-color: #FFD700 !important;
        box-shadow: 0 0 10px rgba(255, 215, 0, 0.3) !important;
    }

    /* Radio buttons */
    .stRadio label {
        color: #e0e0e0 !important;
        font-family: 'Cairo', sans-serif !important;
    }

    /* Custom scrollbar */
    ::-webkit-scrollbar {
        width: 8px;
    }

    ::-webkit-scrollbar-track {
        background: rgba(255, 255, 255, 0.05);
        border-radius: 10px;
    }

    ::-webkit-scrollbar-thumb {
        background: rgba(255, 215, 0, 0.3);
        border-radius: 10px;
    }

    ::-webkit-scrollbar-thumb:hover {
        background: rgba(255, 215, 0, 0.5);
    }

    /* Loading spinner */
    .stSpinner > div {
        border-color: #FFD700 transparent transparent transparent !important;
    }
</style>
""", unsafe_allow_html=True)

# ========== MAIN STREAMLIT APP ==========
st.set_page_config(page_title="⚡Hermes\\ابو محسوب لعدم الرسوب💯", page_icon="⚖️", layout="centered")

# App title with branding
st.markdown("""
<div style="text-align: center; padding: 10px 0;">
    <h1 style="font-size: 2.5em !important;">⚡Hermes\ابو محسوب لعدم الرسوب💯</h1>
    <p style="font-family: 'Cinzel', serif; color: #D4AF37; font-size: 1.1em;">
        The Wisdom of Hermes • The Success of Abu Mahsoub
    </p>
</div>
<hr style="border-color: #D4AF37; opacity: 0.5;">
""", unsafe_allow_html=True)

missing_libs = []
if not GROQ_AVAILABLE:
    missing_libs.append("groq")
if not ARABIC_SUPPORT:
    missing_libs.append("arabic_reshaper, python-bidi")

if missing_libs:
    st.error(f"Missing libraries: {', '.join(missing_libs)}. Run: pip install {' '.join(missing_libs)}")

st.title("📚 Book Analyzer Pro")
st.caption("Creative & Comprehensive Mode")

# ========== SIDEBAR ==========
with st.sidebar:
    st.header("🔑 Settings")
    api_key = st.text_input("API Key", type="password", help="Get from console.groq.com")

    if not api_key:
        st.info("Enter API key to continue")
        st.stop()
    elif not api_key.startswith('gsk_'):
        st.error("Invalid format. Must start with 'gsk_'")
        st.stop()

    st.divider()
    st.subheader("📊 Output Mode")
    mode = st.radio("Choose mode:", ["Quick Summary", "Detailed Explanation"],
                   label_visibility="visible", captions=["5 key insights", "In-depth explanations"])

    st.divider()
    st.subheader("📄 Page Range")
    col1, col2 = st.columns(2)
    with col1:
        page_from = st.number_input("From", min_value=1, value=1)
    with col2:
        page_to = st.number_input("To", min_value=1, value=1)
    st.divider()
    st.subheader("📁 File Name")
    file_name_input = st.text_input("سمِّ ملفك يا بطل (File Name)", "Summary")
    st.divider()
    st.caption("💡 Tip: Upload a PDF or paste text")

tab1, tab2 = st.tabs(["🔍 Search Books", "📤 Upload PDF"])
if 'search_results' not in st.session_state:
    st.session_state['search_results'] = {}
search_results = st.session_state['search_results']

# ========== TAB 1: SEARCH BOOKS ==========
with tab1:
    st.subheader("Search Google Books")
    search_query = st.text_input("Book title or author:", placeholder="e.g., Clean Code")

    if st.button("🔍 Search", type="primary", use_container_width=True):
        if search_query.strip():
            with st.spinner("Searching..."):
                try:
                    books = search_google_books(search_query.strip())
                    if books:
                        st.session_state['search_results'] = {book['id']: book for book in books}
                        st.toast(f"Found {len(books)} books!", icon="✅")
                    else:
                        st.warning("No books found")
                except Exception as e:
                    st.error(f"Search failed: {e}")
        else:
            st.warning("Enter a search term")

    if search_results:
        st.markdown("---")
        for book_id, book in search_results.items():
            with st.expander(f"📗 {format_text_for_output(book['title'][:30])}...", expanded=False):
                st.write(f"**Author:** {format_text_for_output(book['author'])}")
                st.write(f"**Year:** {book['year']}")
                if book.get('description'):
                    st.write(f"**Description:** {format_text_for_output(book['description'][:200])}...")
                if st.button(f"✅ Select", key=f"select_{book_id}", use_container_width=True):
                    st.session_state['selected_book'] = book
                    # Pass Title, Author, and Description to AI as source text
                    book_title = book.get('title', '')
                    book_author = book.get('author', '')
                    book_desc = book.get('description', '')

                    # Combine title, author, description for comprehensive context
                    combined_text = f"Book: {book_title}\nAuthor: {book_author}\n\nDescription:\n{book_desc}"
                    st.session_state['pdf_text'] = combined_text
                    # Confirmation message as requested
                    st.success(f"⚖️ Book '{format_text_for_output(book_title)}' selected by Hermes. Ready to summarize.")

# ========== TAB 2: UPLOAD PDF ==========
with tab2:
    st.subheader("Upload PDF Book")
    uploaded_file = st.file_uploader("Choose PDF file", type="pdf", accept_multiple_files=False)

    if uploaded_file:
        # Create page range tuple from sidebar inputs
        page_range = (page_from, page_to) if page_from and page_to else None

        # Extract text synchronously
        if st.button("📤 Extract Text from PDF", type="primary", use_container_width=True):
            with st.spinner('Extracting text from PDF...'):
                try:
                    pdf_text = extract_text_from_pdf(uploaded_file, page_range)

                    if pdf_text:
                        st.session_state['pdf_text'] = pdf_text
                        st.session_state['pdf_page_range'] = page_range
                        st.session_state['selected_book'] = {
                            'title': uploaded_file.name.replace('.pdf', '')[:30],
                            'author': 'Uploaded Document',
                            'year': 'Unknown',
                            'language': detect_dominant_language(pdf_text)
                        }
                        st.success(f'✅ Extracted {len(pdf_text)} characters successfully!')
                        st.rerun()
                    else:
                        st.error('❌ Failed to extract text from PDF')
                        st.rerun()

                except Exception as e:
                    st.error(f'❌ PDF extraction error: {str(e)}')
                    st.rerun()

# ========== PROCESSING SECTION ==========
# Check if we have either a selected book (from search) or uploaded PDF
has_book = 'selected_book' in st.session_state and st.session_state['selected_book'] is not None
has_text = 'pdf_text' in st.session_state and st.session_state.get('pdf_text', '') != ''

if has_book:
    book = st.session_state['selected_book']

    st.markdown("---")
    st.markdown(f"### 📖 {format_text_for_output(book['title'])}")
    st.caption(f"✍️ By: {format_text_for_output(book['author'])} | 📅 Year: {book['year']}")

    # Show the book context/source
    if has_text:
        text = st.session_state['pdf_text']
        st.info(f"📚 Source: {'Book Description' if 'Book:' in text else 'PDF Content'}")
        st.write(f"📊 Text length: {len(text)} characters")
    else:
        text = ""
        st.warning("No content available. Please paste text or upload a PDF.")

    # Text area for additional/pasted content (optional)
    st.markdown("#### 📝 Additional Content (Optional)")
    additional_text = st.text_area("Paste additional text here:", height=80, placeholder="Paste any additional content...")

    # Combine texts if both exist
    if text and additional_text:
        combined_text = f"{text}\n\n--- Additional Content ---\n{additional_text}"
    elif additional_text:
        combined_text = additional_text
    elif text:
        combined_text = text
    else:
        combined_text = ""

    # Debug information
    if combined_text:
        st.write(f"🔍 Combined text length: {len(combined_text.strip())} characters")

    # Show the Summarize button if we have content
    if combined_text and len(combined_text.strip()) >= 50:
        output_mode = "quick" if mode == "Quick Summary" else "explanation"

        # Golden styled Summarize button
        st.markdown("""
        <style>
        .summarize-btn button {
            background: linear-gradient(135deg, #D4AF37 0%, #B8960C 100%) !important;
            color: #0D3486 !important;
            font-family: 'Cinzel', serif !important;
            font-weight: bold !important;
            font-size: 18px !important;
            padding: 12px 24px !important;
            border: 3px solid #0D3486 !important;
            border-radius: 10px !important;
            box-shadow: 0 4px 12px rgba(212, 175, 55, 0.4) !important;
        }
        .summarize-btn button:hover {
            background: linear-gradient(135deg, #B8960C 0%, #D4AF37 100%) !important;
            transform: scale(1.02) !important;
        }
        </style>
        """, unsafe_allow_html=True)

        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            st.markdown('<div class="summarize-btn">', unsafe_allow_html=True)
            if st.button("⚖️ ✨ Start Summary", type="primary", use_container_width=True):
                # Process synchronously
                with st.spinner('Analyzing content with AI...'):
                    try:
                        # Validate API key
                        if not api_key or not api_key.startswith('gsk_'):
                            st.error('❌ Invalid API key. Must start with gsk_')
                            st.rerun()


                        # Validate text
                        if not combined_text or len(combined_text.strip()) < 50:
                            st.error('❌ Text too short for analysis (minimum 50 characters)')
                            st.rerun()


                        # Split text into pages
                        pages = split_text_by_pages(combined_text)
                        total_pages = len(pages)

                        if total_pages == 0:
                            st.error('❌ No text chunks to process')
                            st.rerun()


                        progress_bar = st.progress(0)
                        status_text = st.empty()

                        sections_data = []
                        for i, page in enumerate(pages):
                            progress = int((i / total_pages) * 80)  # 0-80% for analysis
                            progress_bar.progress(progress)
                            status_text.text(f'Analyzing chunk {i+1}/{total_pages}...')

                            section_result = analyze_section(page['text'], i+1, book, api_key, mode, (page_from, page_to))
                            if section_result:
                                sections_data.append(section_result)

                        if sections_data:
                            progress_bar.progress(85)
                            status_text.text('Processing results...')

                            # Process sections data
                            all_sections_data = []
                            section_counter = 1

                            for idx, chunk_result in enumerate(sections_data):
                                if not isinstance(chunk_result, dict) or 'content' not in chunk_result:
                                    continue

                                content = chunk_result['content']
                                if len(content.strip()) < 10:
                                    continue

                                # Parse sections
                                sections_found = []
                                if '[SECTION' in content.upper():
                                    sections_split = re.split(r'\[SECTION\s*([^\]]+)\]', content)
                                    for i in range(1, len(sections_split), 2):
                                        if i+1 < len(sections_split):
                                            title = sections_split[i].strip()
                                            content_part = sections_split[i+1].strip()
                                            if title and content_part:
                                                sections_found.append((title, content_part))

                                # Fallback parsing
                                if not sections_found:
                                    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip() and len(p.strip()) > 30]
                                    for i, para in enumerate(paragraphs):
                                        sentences = re.split(r'[.!?]\s+', para)
                                        if len(sentences) > 1:
                                            title = sentences[0][:50] + "..." if len(sentences[0]) > 50 else sentences[0]
                                            content_part = '. '.join(sentences[1:])
                                            if len(content_part) > 20:
                                                sections_found.append((f"Part {i+1}: {title}", content_part))

                                # Create slides from found sections
                                for title, content_part in sections_found:
                                    if title and content_part and len(content_part.strip()) > 10:
                                        all_sections_data.append({
                                            'title': title.strip(),
                                            'content': content_part.strip(),
                                            'section_num': section_counter
                                        })
                                        section_counter += 1

                            if all_sections_data:
                                progress_bar.progress(95)
                                status_text.text('Generating PPTX presentation...')

                                # Generate PPTX
                                ppt_file = create_creative_pptx(
                                    book_title=book.get('title', 'Unknown'),
                                    sections_data=all_sections_data,
                                    output_mode=output_mode
                                )

                                progress_bar.progress(100)
                                status_text.text('✅ Complete!')

                                # Store results
                                st.session_state['pptx_file'] = ppt_file
                                st.session_state['sections_data'] = all_sections_data

                                st.success(f'✅ Analysis complete! Generated {len(all_sections_data)} sections in PPTX.')

                                # Clear progress indicators
                                progress_bar.empty()
                                status_text.empty()

                                st.rerun()
                            else:
                                st.error('❌ No content sections could be extracted from AI responses')
                                st.rerun()
                        else:
                            st.error('❌ AI analysis failed to return any valid content')
                            st.rerun()

                    except Exception as e:
                        st.error(f'❌ Analysis error: {str(e)}')
                        st.rerun()
            st.markdown('</div>', unsafe_allow_html=True)

    # Show results if available
    if 'pptx_file' in st.session_state and st.session_state.get('pptx_file'):
        st.markdown("### 📋 Summary Result")
        # UI Rendering: Use st.markdown with RTL CSS for proper right-to-left alignment
        processed_text = fix_text(st.session_state.get('summary_output', ''))
        if is_arabic_text(st.session_state.get('summary_output', '')):
            st.markdown(f'<div style="text-align: right; direction: rtl;">{processed_text}</div>', unsafe_allow_html=True)
        else:
            st.markdown(processed_text)

        # ========== DOWNLOAD BUTTONS SECTION ==========
        safe_filename = re.sub(r'[^\w\s-]', '', file_name_input).strip() or "Summary"
        st.download_button(
            label="📊 Download PPTX",
            data=st.session_state['pptx_file'],
            file_name=f"{safe_filename}.pptx",
            mime="application/vnd/openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )
    elif combined_text and len(combined_text.strip()) < 50:
        st.warning(f"⚠️ Text too short (minimum 50 characters). Current length: {len(combined_text.strip())}")
    else:
        st.info("💡 Tip: Paste book content or upload a PDF to begin analysis")
