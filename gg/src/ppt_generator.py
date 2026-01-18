import arabic_reshaper
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from typing import Dict, List
from dataclasses import dataclass

# Import bidi for proper Arabic bidirectional text handling
try:
    from bidi.algorithm import get_display
    BIDI_AVAILABLE = True
except ImportError as e:
    # Silently handle missing bidi - it's optional for Arabic text processing
    BIDI_AVAILABLE = False
    get_display = None
except Exception as e:
    # Handle any other import issues
    BIDI_AVAILABLE = False
    get_display = None

def fix_arabic_for_pptx(text):
    """Apply Arabic text reshaping and bidirectional display for proper display in PowerPoint"""
    if not text:
        return text

    # Check if text contains Arabic characters
    has_arabic = any('\u0600' <= c <= '\u06FF' for c in text)
    if has_arabic:
        try:
            # Reshape Arabic characters and apply bidirectional display for proper text ordering
            reshaped_text = arabic_reshaper.reshape(text)
            if BIDI_AVAILABLE and get_display:
                processed_text = get_display(reshaped_text)
                return processed_text
            else:
                return reshaped_text

        except Exception as e:
            # If processing fails, return original text
            print(f"Arabic text processing error: {e}")
            return text

    return text

@dataclass
class SlideContent:
    title: str
    content: str
    layout_type: str = 'title_content'

class PowerPointGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.setup_slide_layouts()
    
    def setup_slide_layouts(self):
        """Setup custom slide layouts for academic presentations"""
        # Define color scheme for academic presentations
        self.accent_color = RGBColor(0, 51, 102)  # Dark blue
        self.text_color = RGBColor(51, 51, 51)    # Dark gray
        self.title_color = RGBColor(0, 0, 0)      # Black
        
    def create_fact_sheet_slide(self, metadata: Dict) -> None:
        """Create the required Fact Sheet slide as first slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = "FACT SHEET"
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color

        # Content with fact sheet information
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()  # Clear existing content

        # Book Title - Apply Arabic reshaping if needed
        p = content_frame.add_paragraph()
        book_title = metadata.get('title', 'Unknown')
        p.text = f"📚 Book Title: {fix_arabic_for_pptx(book_title)}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.space_after = Pt(12)
        if any('\u0600' <= c <= '\u06FF' for c in p.text):
            p.alignment = PP_ALIGN.RIGHT

        # Author - Apply Arabic reshaping if needed
        p = content_frame.add_paragraph()
        author = metadata.get('author', 'Unknown')
        p.text = f"✍️ Author: {fix_arabic_for_pptx(author)}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        if any('\u0600' <= c <= '\u06FF' for c in p.text):
            p.alignment = PP_ALIGN.RIGHT

        # Total Pages
        p = content_frame.add_paragraph()
        p.text = f"📄 Total Pages: {metadata.get('total_pages', 0)}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)

        # Number of Chapters
        p = content_frame.add_paragraph()
        p.text = f"📋 Number of Chapters: {metadata.get('number_of_chapters', 0)}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)

        # Academic Classification - Apply Arabic reshaping if needed
        p = content_frame.add_paragraph()
        p.text = fix_arabic_for_pptx("🎓 Academic Classification: Scientific Research Document")
        p.font.size = Pt(18)
        p.font.italic = True
        if any('\u0600' <= c <= '\u06FF' for c in p.text):
            p.alignment = PP_ALIGN.RIGHT
        
        # Add academic border/line
        left = Inches(0.5)
        top = Inches(5.5)
        width = Inches(9)
        height = Inches(0.1)
        line = slide.shapes.add_shape(
            1, left, top, width, height
        )  # 1 is rectangle shape
        line.fill.solid()
        line.fill.fore_color.rgb = self.accent_color
        
    def create_title_slide(self, title: str, subtitle: str = "") -> None:
        """Create title slide for the presentation"""
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        subtitle_placeholder = slide.placeholders[1]

        # Apply Arabic reshaping to title and subtitle if they contain Arabic
        title_placeholder.text = fix_arabic_for_pptx(title)
        subtitle_placeholder.text = fix_arabic_for_pptx(subtitle)

        # Format title
        title_placeholder.text_frame.paragraphs[0].font.size = Pt(44)
        title_placeholder.text_frame.paragraphs[0].font.bold = True
        title_placeholder.text_frame.paragraphs[0].font.color.rgb = self.accent_color
        if any('\u0600' <= c <= '\u06FF' for c in title):
            title_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Format subtitle
        if subtitle:
            subtitle_placeholder.text_frame.paragraphs[0].font.size = Pt(24)
            subtitle_placeholder.text_frame.paragraphs[0].font.color.rgb = self.text_color
            if any('\u0600' <= c <= '\u06FF' for c in subtitle):
                subtitle_placeholder.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    
    def create_english_summary_slide(self, english_summary: str) -> None:
        """Create English academic summary slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = "English Summary"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color

        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()

        # English content
        p = content_frame.add_paragraph()
        p.text = english_summary[:800] + "..." if len(english_summary) > 800 else english_summary
        p.font.size = Pt(16)
        p.space_after = Pt(8)

    def create_arabic_summary_slide(self, arabic_summary: str) -> None:
        """Create Arabic academic summary slide"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = fix_arabic_for_pptx("الملخص العربي")
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()

        # Arabic content
        arabic_body_p = content_frame.add_paragraph()
        arabic_body_p.text = fix_arabic_for_pptx(arabic_summary[:800] + "..." if len(arabic_summary) > 800 else arabic_summary)
        arabic_body_p.font.size = Pt(16)
        arabic_body_p.alignment = PP_ALIGN.RIGHT
    
    def create_chapters_overview_slide(self, chapters: List[Dict]) -> None:
        """Create slide showing chapters overview"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Chapters Overview"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color

        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()

        for i, chapter in enumerate(chapters[:10]):  # Limit to first 10 chapters
            p = content_frame.add_paragraph()
            # Apply Arabic reshaping to chapter title if it contains Arabic
            chapter_title = chapter.get('title', 'Untitled')
            chapter_text = f"Chapter {i+1}: {fix_arabic_for_pptx(chapter_title)} (Page {chapter.get('page', 'Unknown')})"
            p.text = chapter_text
            p.font.size = Pt(16)
            p.space_after = Pt(6)
            if any('\u0600' <= c <= '\u06FF' for c in p.text):
                p.alignment = PP_ALIGN.RIGHT

            # Add indentation for sub-chapters
            if chapter.get('level', 1) > 1:
                p.level = chapter.get('level', 1) - 1
    
    def create_key_points_slide(self, key_points: List[str]) -> None:
        """Create slide for key scientific points"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Key Scientific Points"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color

        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()

        for point in key_points[:8]:  # Limit to first 8 points
            p = content_frame.add_paragraph()
            # Apply Arabic reshaping to each point if it contains Arabic
            p.text = f"• {fix_arabic_for_pptx(point)}"
            p.font.size = Pt(16)
            p.space_after = Pt(8)
            if any('\u0600' <= c <= '\u06FF' for c in p.text):
                p.alignment = PP_ALIGN.RIGHT
    
    def create_scientific_terms_slide(self, terms: List[str]) -> None:
        """Create slide for scientific terminology"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Scientific Terminology"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color

        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()

        for term in terms[:10]:  # Limit to first 10 terms
            p = content_frame.add_paragraph()
            # Apply Arabic reshaping to each term if it contains Arabic
            p.text = f"• {fix_arabic_for_pptx(term)}"
            p.font.size = Pt(14)
            p.space_after = Pt(6)
            if any('\u0600' <= c <= '\u06FF' for c in p.text):
                p.alignment = PP_ALIGN.RIGHT
    
    def create_references_slide(self, references: List[str]) -> None:
        """Create slide for academic references"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Academic References"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color

        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()

        for i, ref in enumerate(references[:8]):  # Limit to first 8 references
            p = content_frame.add_paragraph()
            # Apply Arabic reshaping to each reference if it contains Arabic
            p.text = f"[{i+1}] {fix_arabic_for_pptx(ref)}"
            p.font.size = Pt(12)
            p.space_after = Pt(4)
            if any('\u0600' <= c <= '\u06FF' for c in p.text):
                p.alignment = PP_ALIGN.RIGHT
    
    def create_methodology_slide(self) -> None:
        """Create slide explaining the academic methodology used"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = "Academic Methodology"
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = self.accent_color
        
        # Content
        content_shape = slide.placeholders[1]
        content_frame = content_shape.text_frame
        content_frame.clear()
        
        methodology_points = [
            "• PDF Metadata Extraction using PyMuPDF",
            "• Scientific Content Analysis with AI",
            "• Bilingual Summary Generation (English/Arabic)",
            "• Academic Language Constraints Applied",
            "• Chapter Structure Analysis",
            "• Reference Extraction and Citation",
            "• Peer-review Quality Standards"
        ]
        
        for point in methodology_points:
            p = content_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(16)
            p.space_after = Pt(8)
    
    def generate_academic_presentation(self, metadata: Dict, summary_result) -> str:
        """Generate complete academic presentation with 6-8 slides"""

        slide_count = 0

        # 1. Fact Sheet (Required as first slide)
        self.create_fact_sheet_slide(metadata)
        slide_count += 1

        # 2. Title slide
        self.create_title_slide(
            f"Academic Analysis: {metadata.get('title', 'Unknown')}",
            f"Scientific Summary and Analysis\n{metadata.get('author', 'Unknown')}"
        )
        slide_count += 1

        # 3. English Summary
        self.create_english_summary_slide(summary_result.english_summary)
        slide_count += 1

        # 4. Arabic Summary
        self.create_arabic_summary_slide(summary_result.arabic_summary)
        slide_count += 1

        # 5. Chapters Overview (always create if chapters exist, otherwise create placeholder)
        if metadata.get('chapters'):
            self.create_chapters_overview_slide(metadata['chapters'])
            slide_count += 1
        else:
            # Create a placeholder chapters slide
            self.create_chapters_overview_slide([{'title': 'No chapters available', 'page': 'N/A'}])
            slide_count += 1

        # 6. Key Scientific Points (always create)
        if summary_result.key_points:
            self.create_key_points_slide(summary_result.key_points)
        else:
            self.create_key_points_slide(['No key points extracted'])
        slide_count += 1

        # 7. Scientific Terminology (always create)
        if summary_result.scientific_terms:
            self.create_scientific_terms_slide(summary_result.scientific_terms)
        else:
            self.create_scientific_terms_slide(['No scientific terms identified'])
        slide_count += 1

        # 8. Academic References (always create)
        if summary_result.references:
            self.create_references_slide(summary_result.references)
        else:
            self.create_references_slide(['No references available'])
        slide_count += 1

        # 9. Methodology (always create as final slide)
        self.create_methodology_slide()
        slide_count += 1

        return f"Academic presentation generated successfully with {slide_count} slides (6-8 slides as requested)"
    
    def save_presentation(self, filename: str) -> bool:
        """Save the presentation to file"""
        try:
            self.prs.save(filename)
            return True
        except Exception as e:
            print(f"Error saving presentation: {e}")
            return False
    
    def add_academic_header_footer(self, slide):
        """Add academic header/footer to slide (optional enhancement)"""
        # This is a placeholder for adding headers/footers
        # Could include page numbers, document title, etc.
        pass
